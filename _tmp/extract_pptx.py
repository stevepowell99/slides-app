"""One-off: dump structured content of EES2025 coffee break.pptx."""
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

PPTX = Path(r"C:\Users\Zoom\My Drive (hello@causalmap.app)\Causal Map\10-19 Outreach - marketing - presentations - academic - theory\13 Presentations\!individual presentations\EES2025 coffee break.pptx")
OUT = Path(__file__).parent / "extracted.json"

p = Presentation(PPTX)

def shape_to_dict(shape, slide_w, slide_h):
    d = {
        "name": shape.name,
        "type": shape.shape_type.__repr__() if shape.shape_type else None,
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }
    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            runs = [{"text": r.text, "bold": r.font.bold, "italic": r.font.italic, "size": r.font.size and r.font.size.pt} for r in para.runs]
            paragraphs.append({"level": para.level, "runs": runs, "text": para.text})
        d["paragraphs"] = paragraphs
    if shape.shape_type and "PICTURE" in str(shape.shape_type):
        d["is_picture"] = True
    return d

slides = []
for i, slide in enumerate(p.slides, 1):
    title = ""
    if slide.shapes.title:
        title = slide.shapes.title.text
    shapes = [shape_to_dict(s, p.slide_width, p.slide_height) for s in slide.shapes]
    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
    slides.append({"index": i, "title": title, "shapes": shapes, "notes": notes})

OUT.write_text(json.dumps({
    "slide_width": p.slide_width,
    "slide_height": p.slide_height,
    "slides": slides
}, indent=2), encoding="utf-8")

print(f"Wrote {OUT}")
for s in slides:
    print(f"  slide {s['index']}: {s['title']!r}")
